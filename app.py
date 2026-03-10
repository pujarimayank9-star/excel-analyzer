# app.py

import streamlit as st
import pandas as pd
import matplotlib.pyplot as plt
import random
from pptx import Presentation
from pptx.util import Inches
import tempfile

st.set_page_config(layout="wide")
st.title("Universal Excel Data Analyst")

file = st.file_uploader("Upload Excel / CSV", type=["xlsx","csv"])

if file:

    # ---------------- LOAD DATA ----------------
    if file.name.endswith(".csv"):
        df = pd.read_csv(file)
    else:
        df = pd.read_excel(file)

    st.header("Dataset Overview")
    st.write(f"Rows: {df.shape[0]} | Columns: {df.shape[1]}")
    st.dataframe(df.head())

    numeric_cols = df.select_dtypes(include=["int64","float64"]).columns
    cat_cols = df.select_dtypes(include=["object"]).columns

    insights = []
    recommendations = []
    charts = []

    # ---------------- INSIGHT FUNCTIONS ----------------

    def histogram_insight(col, avg, mn, mx):
        t = [
            f"{col} ranges between {mn} and {mx} with an average of {round(avg,2)} indicating the spread of values.",
            f"The distribution of {col} spans from {mn} to {mx} with mean around {round(avg,2)}.",
            f"{col} values fluctuate between {mn} and {mx}, suggesting variability in the dataset."
        ]
        return random.choice(t)

    def bar_insight(col, data):
        best = data.idxmax()
        worst = data.idxmin()

        t = [
            f"In {col}, {best} shows the highest presence while {worst} appears least.",
            f"{best} dominates the {col} comparison whereas {worst} contributes the lowest.",
            f"{best} leads performance in {col} while {worst} has minimal representation."
        ]

        return random.choice(t), best, worst

    def pie_insight(col, counts):
        top = counts.idxmax()
        percent = round((counts.max()/counts.sum())*100,1)

        t = [
            f"{top} holds the largest share in {col} contributing about {percent}%.",
            f"{top} dominates the distribution of {col} with roughly {percent}% share.",
            f"{top} accounts for the highest proportion of {col}."
        ]

        return random.choice(t)

    def recommendation(best, worst, col):
        t = [
            f"Investigate why {best} performs strongly in {col} and replicate similar strategies.",
            f"Focus on improving the performance of {worst} within {col}.",
            f"Practices used by {best} could help improve other segments in {col}."
        ]
        return random.choice(t)

    # ---------------- NUMERIC ANALYSIS ----------------

    st.header("Numeric Analysis")

    for col in numeric_cols:

        st.subheader(col)

        avg = df[col].mean()
        mx = df[col].max()
        mn = df[col].min()

        c1,c2,c3 = st.columns(3)
        c1.metric("Average", round(avg,2))
        c2.metric("Max", mx)
        c3.metric("Min", mn)

        fig, ax = plt.subplots()
        df[col].hist(ax=ax)

        st.pyplot(fig)

        insight = histogram_insight(col, avg, mn, mx)
        st.write("Insight:", insight)

        insights.append(insight)
        charts.append((col,fig))

    # ---------------- CATEGORY ANALYSIS ----------------

    st.header("Category Analysis")

    for cat in cat_cols:

        if df[cat].nunique() <= 20:

            st.subheader(cat)

            counts = df[cat].value_counts()

            # BAR
            fig1, ax = plt.subplots()
            counts.plot(kind="bar", ax=ax)

            st.pyplot(fig1)

            insight, best, worst = bar_insight(cat, counts)
            st.write("Insight:", insight)

            rec = recommendation(best, worst, cat)
            st.write("Recommendation:", rec)

            insights.append(insight)
            recommendations.append(rec)

            charts.append((cat+"_bar", fig1))

            # PIE
            fig2, ax = plt.subplots()
            counts.plot(kind="pie", autopct="%1.1f%%", ax=ax)

            st.pyplot(fig2)

            pie_i = pie_insight(cat, counts)
            st.write("Insight:", pie_i)

            insights.append(pie_i)

            charts.append((cat+"_pie", fig2))

    # ---------------- CORRELATION ANALYSIS ----------------

    if len(numeric_cols) > 1:

        st.header("Correlation Analysis")

        corr = df[numeric_cols].corr()

        fig, ax = plt.subplots(figsize=(6,5))
        cax = ax.imshow(corr, cmap="coolwarm")

        ax.set_xticks(range(len(corr.columns)))
        ax.set_yticks(range(len(corr.columns)))

        ax.set_xticklabels(corr.columns, rotation=45)
        ax.set_yticklabels(corr.columns)

        plt.colorbar(cax)

        for i in range(len(corr.columns)):
            for j in range(len(corr.columns)):
                ax.text(j, i, round(corr.iloc[i,j],2),
                        ha="center", va="center", color="black")

        st.pyplot(fig)

        charts.append(("correlation", fig))

        st.subheader("Correlation Insights")

        for i in range(len(corr.columns)):
            for j in range(i+1, len(corr.columns)):

                col1 = corr.columns[i]
                col2 = corr.columns[j]
                value = corr.iloc[i,j]

                if abs(value) < 0.2:
                    relation = "very weak relationship"
                elif abs(value) < 0.4:
                    relation = "weak relationship"
                elif abs(value) < 0.7:
                    relation = "moderate relationship"
                else:
                    relation = "strong relationship"

                if value > 0:
                    direction = "positive"
                else:
                    direction = "negative"

                insight = (
                    f"{col1} and {col2} show a {relation} ({direction}) "
                    f"with correlation value {round(value,2)}. "
                    f"This means changes in {col1} are associated with similar "
                    f"changes in {col2}." if value>0 else
                    f"{col1} and {col2} show a {relation} ({direction}) "
                    f"with correlation value {round(value,2)}. "
                    f"As {col1} increases, {col2} tends to decrease."
                )

                st.write("Insight:", insight)
                insights.append(insight)

                rec = (
                    f"Monitor interaction between {col1} and {col2} as their relationship "
                    f"may influence performance trends."
                )

                st.write("Recommendation:", rec)
                recommendations.append(rec)

    # ---------------- PPT GENERATION ----------------

    if st.button("Generate PPT Report"):

        prs = Presentation()

        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Automated Data Analysis Report"
        slide.placeholders[1].text = "Generated from Excel Dataset"

        i = 0

        while i < len(charts):

            name, chart1 = charts[i]

            tmp1 = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
            chart1.savefig(tmp1.name)

            chart2 = None

            if i+1 < len(charts) and "pie" in charts[i+1][0]:
                chart2 = charts[i+1][1]

            slide = prs.slides.add_slide(prs.slide_layouts[5])
            slide.shapes.title.text = name.replace("_bar","").replace("_pie","")

            slide.shapes.add_picture(
                tmp1.name,
                Inches(0.5),
                Inches(1.5),
                height=Inches(3)
            )

            if chart2:

                tmp2 = tempfile.NamedTemporaryFile(delete=False, suffix=".png")
                chart2.savefig(tmp2.name)

                slide.shapes.add_picture(
                    tmp2.name,
                    Inches(5),
                    Inches(1.5),
                    height=Inches(3)
                )

                i += 1

            if i < len(insights):

                textbox = slide.shapes.add_textbox(
                    Inches(0.5),
                    Inches(4.8),
                    Inches(9),
                    Inches(1)
                )

                tf = textbox.text_frame
                tf.text = "Insight: " + insights[i]

            if i < len(recommendations):

                textbox2 = slide.shapes.add_textbox(
                    Inches(0.5),
                    Inches(5.6),
                    Inches(9),
                    Inches(1)
                )

                tf2 = textbox2.text_frame
                tf2.text = "Recommendation: " + recommendations[i]

            i += 1

        prs.save("analysis_report.pptx")

        with open("analysis_report.pptx","rb") as f:

            st.download_button(
                "Download PPT",
                f,
                file_name="analysis_report.pptx"
            )
