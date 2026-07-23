from pptx import Presentation
from pptx.util import Inches
import streamlit as st
import tempfile
import os

def ppt_download(df):

    st.header("📄 PowerPoint Report")

    if st.button("Generate PPT"):

        prs = Presentation()

        # ---------- TITLE ----------
        slide = prs.slides.add_slide(prs.slide_layouts[0])

        slide.shapes.title.text = "AI Smart Excel Analysis Report"

        slide.placeholders[1].text = (
            f"Rows : {df.shape[0]}\n"
            f"Columns : {df.shape[1]}"
        )

        # ---------- DATA PREVIEW ----------
        slide = prs.slides.add_slide(prs.slide_layouts[5])

        slide.shapes.title.text = "Dataset Preview"

        rows = min(10, len(df))

        cols = min(6, len(df.columns))

        table = slide.shapes.add_table(
            rows + 1,
            cols,
            Inches(0.3),
            Inches(0.8),
            Inches(9),
            Inches(3.5)
        ).table

        for c in range(cols):
            table.cell(0, c).text = str(df.columns[c])

        for r in range(rows):
            for c in range(cols):
                table.cell(r + 1, c).text = str(df.iloc[r, c])

        # ---------- SUMMARY ----------
        slide = prs.slides.add_slide(prs.slide_layouts[1])

        slide.shapes.title.text = "Summary"

        tf = slide.placeholders[1].text_frame

        tf.text = f"Rows : {df.shape[0]}"

        tf.add_paragraph().text = f"Columns : {df.shape[1]}"

        tf.add_paragraph().text = f"Missing Values : {df.isnull().sum().sum()}"

        tf.add_paragraph().text = f"Duplicate Rows : {df.duplicated().sum()}"

        filename = "AI_Analysis_Report.pptx"

        prs.save(filename)

        with open(filename, "rb") as f:

            st.download_button(
                "⬇ Download PPT",
                f,
                file_name=filename
            )

        os.remove(filename)
